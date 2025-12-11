import math
import json
import os
import pandas as pd
import psycopg2
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_TICK_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Cm
from datetime import datetime
from dateutil.relativedelta import relativedelta

TEMPLATE_CANDIDATES = [
    "ppt_template.pptx",
    "ppt-template.pptx",
]

def locate_ppt_template():
    """Return absolute path to the PPT template, trying known variants."""
    base_dir = os.path.dirname(os.path.abspath(__file__))
    for candidate in TEMPLATE_CANDIDATES:
        template_path = os.path.join(base_dir, candidate)
        if os.path.exists(template_path):
            return template_path
    raise FileNotFoundError(
        "PowerPoint template not found. Expected one of: "
        + ", ".join(TEMPLATE_CANDIDATES)
    )


def safe_int(value, default=0):
    """Convert a value to int while handling None/NaN gracefully."""
    if value is None:
        return default
    if isinstance(value, float) and math.isnan(value):
        return default
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def get_db_connection_for_ppt(username, password):
    """Establishes a connection to the PostgreSQL database for PPT generation."""
    try:
        conn = psycopg2.connect(
            dbname='AutomationDB',
            user=username,
            password=password,
            host='10.193.131.151',
            port='5432'
        )
        return conn
    except psycopg2.Error as e:
        print(f"Error connecting to PostgreSQL: {e}")
        return None

def fetch_data(conn, customer_name, month_year):
    """Fetches data from the database for a specific customer and month."""
    
    end_date = datetime.strptime(month_year, '%Y-%m-%d').date()

    # Get the number of months for historical data from the customer_mapping_table
    customer_mapping_sql = "SELECT no_of_months FROM customer_mapping_table WHERE customer_name = %s AND month_year = %s"
    try:
        customer_mapping_df = pd.read_sql(customer_mapping_sql, conn, params=(customer_name, end_date))
        if customer_mapping_df.empty:
            no_of_months = 6
        else:
            no_of_months = customer_mapping_df['no_of_months'].iloc[0]
    except (IndexError, pd.errors.DatabaseError) as e:
        print(f"Could not determine number of months, defaulting to 6. Error: {e}")
        no_of_months = 6

    # Calculate the start date for fetching historical data
    start_date = end_date - relativedelta(months=no_of_months - 1)
    
    # Fetch data from both tables for the specified period
    customer_mapping_sql = "SELECT * FROM customer_mapping_table WHERE customer_name = %s AND month_year BETWEEN %s AND %s"
    final_computed_sql = "SELECT * FROM final_computed_table WHERE customer_name = %s AND month_year BETWEEN %s AND %s"

    customer_mapping_df = pd.read_sql(customer_mapping_sql, conn, params=(customer_name, start_date, end_date))
    final_computed_df = pd.read_sql(final_computed_sql, conn, params=(customer_name, start_date, end_date))

    if customer_mapping_df.empty or final_computed_df.empty:
        raise ValueError(f"No data found for customer '{customer_name}' between {start_date} and {end_date}.")
    return customer_mapping_df, final_computed_df


def prepare_data_dictionary(customer_mapping_df, final_computed_df, month_year):
    """Prepares the data in a dictionary format similar to the original JSON."""
    
    # Filter data for the specified month
    current_month_date = datetime.strptime(month_year, '%Y-%m-%d').date()
    current_customer_df = customer_mapping_df[customer_mapping_df['month_year'] == current_month_date]
    current_computed_df = final_computed_df[final_computed_df['month_year'] == current_month_date]

    if current_customer_df.empty or current_computed_df.empty:
        raise ValueError(f"Data for the current month ({month_year}) is missing.")
    current_customer_data, current_computed_data = current_customer_df.iloc[0], current_computed_df.iloc[0]

    # --- Global Data ---
    indicator_colors = current_customer_data['indicator_color_code_rules']
    circle_colors = current_customer_data['circle_color_code_rules']

    # --- Slide 1 Data ---
    slide1_data = {
        "Customer_Name": current_customer_data['customer_full_name'],
        "Month": datetime.strptime(month_year, '%Y-%m-%d').strftime('%B %Y'),
        "CSM_Name": current_customer_data['csm_primary']
    }

    # --- Slide 2 Data ---
    availability_chart_data = {
        "Months": [d.strftime('%b-%y') for d in sorted(final_computed_df['month_year'].unique())],
        "Availability": [val * 100 for val in final_computed_df.sort_values(by='month_year')['updated_availability']],
        "SLA": [val * 100 for val in final_computed_df.sort_values(by='month_year')['updated_target']]
    }
    
    slide2_data = {
        "Colour_Rules": current_customer_data['color_map_thresholds_availability'],
        "Indicator": indicator_colors,
        "Circle_Color": circle_colors,
        "Actual_Value": f"{current_computed_data['updated_availability'] * 100:.2f}%",
        "Target_Value": f"{current_computed_data['updated_target'] * 100:.2f}%",
        "Production_Availability_Chart": availability_chart_data,
        "Notes_User_Input": current_customer_data['notes_availability']
    }
    
    # --- Slide 3 Data ---
    prod_licenses = safe_int(current_computed_data['updated_prod_limit'])
    prod_used = safe_int(current_computed_data['updated_prod_used'])
    prod_remaining = prod_licenses - prod_used
    prod_used_percent = round((prod_used * 100) / prod_licenses) if prod_licenses else 0

    test_licenses = safe_int(current_computed_data['updated_test_limit'])
    test_used = safe_int(current_computed_data['updated_test_used'])
    test_remaining = test_licenses - test_used
    test_used_percent = round((test_used * 100) / test_licenses) if test_licenses else 0

    user_license_rows = [
        ["Prod", prod_licenses, prod_used, prod_remaining, prod_used_percent],
        ["Test", test_licenses, test_used, test_remaining, test_used_percent],
    ]

    if current_customer_data['no_of_environments'] == 3:
        dev_licenses = safe_int(current_computed_data['updated_dev_limit'])
        dev_used = safe_int(current_computed_data['updated_dev_used'])
        dev_remaining = dev_licenses - dev_used
        dev_used_percent = round((dev_used * 100) / dev_licenses) if dev_licenses else 0
        user_license_rows.append(["Dev", dev_licenses, dev_used, dev_remaining, dev_used_percent])

    user_counts_chart_data = {
        "Months": [d.strftime('%b-%y') for d in sorted(final_computed_df['month_year'].unique())],
        "Prod": list(final_computed_df.sort_values(by='month_year')['updated_prod_used']),
        "Test": list(final_computed_df.sort_values(by='month_year')['updated_test_used']),
        #"Licenses Available": list(final_computed_df.sort_values(by='month_year')['updated_prod_limit'])
    }
    
    if current_customer_data['no_of_environments'] == 3:
        user_counts_chart_data["Dev"] = list(final_computed_df.sort_values(by='month_year')['updated_dev_used'])
        #user_counts_chart_data["Licenses Available"] = list(final_computed_df.sort_values(by='month_year')['updated_prod_limit'])
    
    user_counts_chart_data["Licenses Available"] = list(final_computed_df.sort_values(by='month_year')['updated_prod_limit'])
    
    slide3_data = {
        "User_License_Utilization_Table": {
            "headers": ["", "Licenses", "Count", "Remaining", "%Used"],
            "rows": user_license_rows
        },
        "Colour_Rules": current_customer_data['color_map_thresholds_users'],
        "Indicator": indicator_colors,
        "Circle_Color": circle_colors,
        "Production_User_Counts_Chart": user_counts_chart_data,
        "Notes_User_Input": current_customer_data['notes_users'],
        "env_count": current_customer_data['no_of_environments']
    }

    # --- Slide 4 Data ---
    prod_storage_used = safe_int(current_computed_data['updated_prod_storage_gb'])
    prod_storage_contract = safe_int(current_computed_data['updated_prod_target_storage_gb'])
    prod_storage_free = prod_storage_contract - prod_storage_used
    prod_storage_used_percent = round((prod_storage_used * 100) / prod_storage_contract, 1) if prod_storage_contract else 0
    prod_storage_free_percent = round((prod_storage_free * 100) / prod_storage_contract, 1) if prod_storage_contract else 0

    test_storage_used = safe_int(current_computed_data['updated_test_storage_gb'])
    test_storage_contract = safe_int(current_computed_data['updated_test_target_storage_gb'])
    test_storage_free = test_storage_contract - test_storage_used
    test_storage_used_percent = round((test_storage_used * 100) / test_storage_contract, 1) if test_storage_contract else 0
    test_storage_free_percent = round((test_storage_free * 100) / test_storage_contract, 1) if test_storage_contract else 0
    
    storage_utilization_rows = [
        ["Prod(GB)", prod_storage_used, prod_storage_contract, prod_storage_free, prod_storage_used_percent, prod_storage_free_percent],
        ["Test(GB)", test_storage_used, test_storage_contract, test_storage_free, test_storage_used_percent, test_storage_free_percent],
    ]

    if current_customer_data['no_of_environments'] == 3:
        dev_storage_used = safe_int(current_computed_data['updated_dev_storage_gb'])
        dev_storage_contract = safe_int(current_computed_data['updated_dev_target_storage_gb'])
        dev_storage_free = dev_storage_contract - dev_storage_used
        dev_storage_used_percent = round((dev_storage_used * 100) / dev_storage_contract, 1) if dev_storage_contract else 0
        dev_storage_free_percent = round((dev_storage_free * 100) / dev_storage_contract, 1) if dev_storage_contract else 0
        storage_utilization_rows.append(["Dev(GB)", dev_storage_used, dev_storage_contract, dev_storage_free, dev_storage_used_percent, dev_storage_free_percent])

    storage_usage_chart_data = {
        "Months": [d.strftime('%b-%y') for d in sorted(final_computed_df['month_year'].unique())],
        "Prod (GB)": list(final_computed_df.sort_values(by='month_year')['updated_prod_storage_gb']),
        "Contracted Maximum": list(final_computed_df.sort_values(by='month_year')['updated_prod_target_storage_gb']),
    }

    slide4_data = {
        "Storage_Utilization_Table": {
            "headers": ["", "Used", "Contract", "Free", "%Used", "%Free"],
            "rows": storage_utilization_rows
        },
        "Colour_Rules": current_customer_data['color_map_thresholds_storage'],
        "Indicator": indicator_colors,
        "Circle_Color": circle_colors,
        "Production_Storage_Usage_Chart": storage_usage_chart_data,
        "Notes_User_Input": current_customer_data['notes_storage']
    }
    
    # --- Slide 5 Data ---
    from dateutil.relativedelta import relativedelta

    current_month = current_computed_data['month_year']
    previous_month = current_month - relativedelta(months=1)

# Get the backlog of the previous month (if it exists)
    prev_row = final_computed_df[final_computed_df['month_year'] == previous_month]

    backlog_active_prev = prev_row['updated_tickets_backlog'].iloc[0] if not prev_row.empty else 0
    case_status_rows = [
        ["Backlog (Active previous months)", backlog_active_prev],
        ["Opened this month", current_computed_data['updated_current_opened_tickets']],
        ["Closed this month", current_computed_data['updated_current_closed_tickets']],
        ["In progress at end of month", current_computed_data['updated_tickets_backlog']]
    ]
    
    case_trend_chart_data = {
        "Months": [d.strftime('%b-%y') for d in sorted(final_computed_df['month_year'].unique())],
        "Opened": list(final_computed_df.sort_values(by='month_year')['updated_tickets_opened']),
        "Closed": list(final_computed_df.sort_values(by='month_year')['updated_tickets_closed']),
        "Open at EOM": list(final_computed_df.sort_values(by='month_year')['updated_tickets_backlog'])
    }
    
    slide5_data = {
        "Case_Status_Table": {
            "headers": ["Status", "Cases"],
            "rows": case_status_rows,
        },
        "Case_Trend_Chart": case_trend_chart_data,
        "Open_Cases_Value": current_computed_data['updated_tickets_backlog']
    }
    
    # --- Slide 7 Data ---
    slide7_data = {
        "Production_User_Counts_Chart": user_counts_chart_data,
        "Production_Availability_Chart": availability_chart_data,
        "Production_Storage_Usage_Chart": storage_usage_chart_data
    }

    return {
        "slide1": slide1_data,
        "slide2": slide2_data,
        "slide3": slide3_data,
        "slide4": slide4_data,
        "slide5": slide5_data,
        "slide7": slide7_data
    }


def delete_table_row(table, row_idx: int):
    """Deletes a row from a table."""
    tbl = table._tbl
    tr = tbl.tr_lst[row_idx]
    tbl.remove(tr)


def generate_presentation(data, output_filename):
    """Generates the PowerPoint presentation with the provided data."""
    template_path = locate_ppt_template()
    prs = Presentation(template_path)
    
    # ---Slide 1---
    slide1_data = data["slide1"]
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.name == "Customer_Name" and shape.has_text_frame:
            p = shape.text_frame.paragraphs[0]
            p.text = slide1_data["Customer_Name"]
            run = p.runs[0]
            font = run.font
            font.size = Pt(40)
            font.color.rgb = RGBColor(255, 255, 255)
        elif shape.name == "Month" and shape.has_text_frame:
            p = shape.text_frame.paragraphs[0]
            p.text = slide1_data["Month"]
            run = p.runs[0]
            font = run.font
            font.size = Pt(18)
            font.color.rgb = RGBColor(255, 255, 255)
        elif shape.name == "CSM_Name" and shape.has_text_frame:
            p = shape.text_frame.paragraphs[0]
            p.text = slide1_data["CSM_Name"]
            run = p.runs[0]
            font = run.font
            font.size = Pt(18)
            font.color.rgb = RGBColor(255, 255, 255)

    # ---Slide 2---
    slide2_data = data["slide2"]
    slide = prs.slides[1]
    color_rules = slide2_data["Colour_Rules"]
    circle_colors = slide2_data["Circle_Color"]
    indicator_colors = slide2_data["Indicator"]
    actual_val = float(slide2_data["Actual_Value"].replace('%', ''))
    
    if actual_val >= color_rules["Color1"]:
        color_key = "Color1"
    elif actual_val >= color_rules["Color2"]:
        color_key = "Color2"
    elif actual_val >= color_rules["Color3"]:
        color_key = "Color3"
    else:
        color_key = "Invalid"
    circle_rgb = circle_colors[color_key]
    indicator_rgb = indicator_colors[color_key]

    for shape in slide.shapes:
        if shape.name == "Target_Value" and shape.has_text_frame:
            shape.text_frame.text = slide2_data["Target_Value"]
        elif shape.name == "Actual_Value" and shape.has_text_frame:
            shape.text_frame.text = slide2_data["Actual_Value"]
        elif shape.name == "Circle_Color":
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*circle_rgb)
        elif shape.name == "Indicator":
            fill = shape.fill
            if fill.type == 3:
                stops = fill.gradient_stops
                if len(stops) >= 2:
                    stops[0].color.rgb = RGBColor(*indicator_rgb)
                    stops[1].color.rgb = RGBColor(255, 255, 255)
        elif shape.name == "Production_Availability_Chart" and shape.has_chart:
            chart = shape.chart
            availability = [val / 100 for val in slide2_data["Production_Availability_Chart"]["Availability"]]
            sla = [val / 100 for val in slide2_data["Production_Availability_Chart"]["SLA"]]
            chart_data = CategoryChartData()
            chart_data.categories = slide2_data["Production_Availability_Chart"]["Months"]
            chart_data.add_series("Availability", availability)
            chart_data.add_series("SLA", sla)
            chart.replace_data(chart_data)
            
            value_axis = chart.value_axis
            value_axis.minimum_scale = 0.93
            value_axis.maximum_scale = 1.0
            value_axis.tick_labels.number_format = '0.00%'

            for series in chart.series:
                if series.name == "Availability":
                    series.has_data_labels = True
                    series.data_labels.number_format = '0.00%'
                    series.data_labels.show_value = True
                else:
                    series.has_data_labels = False
        elif shape.name == "Notes_User_Input" and shape.has_text_frame:
            # ---- Replace existing Notes_User_Input parsing with this ----
            notes_data = slide2_data.get("Notes_User_Input")

            # If DB returned a dict (jsonb -> Python dict), map it to the expected keys
            if isinstance(notes_data, dict):
                note_map = {
                    "Color1": notes_data.get("color1") or notes_data.get("Color1") or "",
                    "Color2": notes_data.get("color2") or notes_data.get("Color2") or "",
                    "Color3": notes_data.get("color3") or notes_data.get("Color3") or "",
                    "Invalid": notes_data.get("invalid") or notes_data.get("Invalid") or ""
                }
            # else:
            #     # existing/robust fallback: try to parse string/list as before
            #     try:
            #         if isinstance(notes_data, str):
            #             parsed = json.loads(notes_data)
            #             if isinstance(parsed, dict):
            #                 note_map = {
            #                     "Color1": parsed.get("color1") or parsed.get("Color1") or "",
            #                     "Color2": parsed.get("color2") or parsed.get("Color2") or "",
            #                     "Color3": parsed.get("color3") or parsed.get("Color3") or "",
            #                     "Invalid": parsed.get("invalid") or parsed.get("Invalid") or ""
            #                 }
            #             elif isinstance(parsed, list):
            #                 note_list = parsed
            #                 note_map = {
            #                     "Color1": note_list[0] if len(note_list) > 0 else "",
            #                     "Color2": note_list[1] if len(note_list) > 1 else "",
            #                     "Color3": note_list[2] if len(note_list) > 2 else "",
            #                     "Invalid": note_list[3] if len(note_list) > 3 else ""
            #                 }
            #             else:
            #                 note_map = {"Color1": str(parsed), "Color2":"", "Color3":"", "Invalid":""}
            #         elif isinstance(notes_data, list):
            #             note_list = notes_data
            #             note_map = {
            #                 "Color1": note_list[0] if len(note_list) > 0 else "",
            #                 "Color2": note_list[1] if len(note_list) > 1 else "",
            #                 "Color3": note_list[2] if len(note_list) > 2 else "",
            #                 "Invalid": note_list[3] if len(note_list) > 3 else ""
            #             }
            #         else:
            #             # fallback: convert to string and put in Color1
            #             note_map = {"Color1": str(notes_data or ""), "Color2":"", "Color3":"", "Invalid":""}
            #     except Exception:
            #         note_map = {"Color1":"","Color2":"","Color3":"","Invalid":""}

            note_value = (note_map.get(color_key) or "").strip()

            if not note_value:
                continue

            # convert escaped newlines to actual newlines and set paragraphs (your existing logic)
            note_value = note_value.replace("\\n", "\n")
            lines = [line.strip() for line in note_value.split("\n") if line.strip()]

            shape.text_frame.text = ""
            shape.text_frame.word_wrap = True

            for i, line in enumerate(lines):
                p = shape.text_frame.paragraphs[0] if i == 0 else shape.text_frame.add_paragraph()
                p.text = line
                run = p.runs[0] if p.runs else p.add_run()
                font = run.font
                font.size = Pt(16)
                font.color.rgb = RGBColor(0, 0, 0)


    # ---Slide 3---
    slide3_data = data["slide3"]
    slide = prs.slides[2]
    env_count = slide3_data.get("env_count")
    headers = slide3_data["User_License_Utilization_Table"]["headers"]
    rows = slide3_data["User_License_Utilization_Table"]["rows"]
    used_val = float(rows[0][4])
    used_test = float(rows[1][4])
    rules = slide3_data["Colour_Rules"]
    
    if used_val >= rules["Color3"] or used_test >= rules["Color3"]:
        color_key = "Color3"
    elif used_val >= rules["Color2"] or used_test >= rules["Color2"]:
        color_key = "Color2"
    elif used_val >= rules["Color1"] or used_test >= rules["Color1"]:
        color_key = "Color1"
    else:
        color_key = "Invalid"
    circle_rgb = slide3_data["Circle_Color"][color_key]
    indicator_rgb = slide3_data["Indicator"][color_key]

    #circle_rgb = tuple(int(x) for x in circle_rgb)
    #is_red = (circle_rgb[0] >= 200 and circle_rgb[1] <= 80 and circle_rgb[2] <= 80)
    #is_red = (color_key == "Color3")
    
    # Determine "red" status for Prod and Test individually
    is_prod_red = used_val >= rules["Color3"]
    is_test_red = used_test >= rules["Color3"]

    shapes_to_remove = []

    for shape in slide.shapes:
        if shape.name == "User_License_Utilization_Table" and shape.has_table:
            table = shape.table

            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
            
            for row_idx, row_data in enumerate(rows, start = 1):
                for col_idx, cell_data in enumerate(row_data):
                    cell = table.cell(row_idx, col_idx)
                    if col_idx == len(row_data) - 1:
                        cell.text = f"{cell_data}%"
                    else:
                        cell.text = str(cell_data)
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(12)
                    p.alignment = PP_ALIGN.RIGHT
            
            total_rows_in_table = len(table.rows)
            needed_data_rows = len(rows)
            rows_to_delete = total_rows_in_table - 1 - needed_data_rows
            if rows_to_delete > 0:
                for _ in range(rows_to_delete):
                    delete_table_row(table, len(table.rows) - 1)
        
        elif shape.name == "Production_User_Counts_Chart" and shape.has_chart:
            chart = shape.chart
            chart_data = CategoryChartData()
            chart_data.categories = slide3_data["Production_User_Counts_Chart"]["Months"]
            if env_count == 2:
                plot = chart.plots[0]
                series_list = list(plot.series)
                for idx, s in enumerate(series_list):
                    if s.name == "Dev":
                         # Remove via XML
                            ser_element = plot._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/chart}ser')[idx]
                            plot._element.remove(ser_element)
                            break
            
            chart_data.categories = slide3_data["Production_User_Counts_Chart"]["Months"]
            
            for series_name, values in slide3_data["Production_User_Counts_Chart"].items():
                if series_name != "Months":
                    chart_data.add_series(series_name, values)
            chart.replace_data(chart_data)
        
        elif shape.name == "Circle_Color":
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*circle_rgb)

        elif shape.name == "Indicator":
            fill = shape.fill
            if fill.type == 3:
                stops = fill.gradient_stops
                if len(stops) >= 2:
                    stops[0].color.rgb = RGBColor(*indicator_rgb)
                    stops[1].color.rgb = RGBColor(255, 255, 255)

        elif shape.name == "Notes_User_Input" and shape.has_text_frame:
                        # ---- Replace existing Notes_User_Input parsing with this ----
            notes_data = slide3_data.get("Notes_User_Input")

            # If DB returned a dict (jsonb -> Python dict), map it to the expected keys
            if isinstance(notes_data, dict):
                note_map = {
                    "Color1": notes_data.get("color1") or notes_data.get("Color1") or "",
                    "Color2": notes_data.get("color2") or notes_data.get("Color2") or "",
                    "Color3": notes_data.get("color3") or notes_data.get("Color3") or "",
                    "Invalid": notes_data.get("invalid") or notes_data.get("Invalid") or ""
                }
            # else:
            #     # existing/robust fallback: try to parse string/list as before
            #     try:
            #         if isinstance(notes_data, str):
            #             parsed = json.loads(notes_data)
            #             if isinstance(parsed, dict):
            #                 note_map = {
            #                     "Color1": parsed.get("color1") or parsed.get("Color1") or "",
            #                     "Color2": parsed.get("color2") or parsed.get("Color2") or "",
            #                     "Color3": parsed.get("color3") or parsed.get("Color3") or "",
            #                     "Invalid": parsed.get("invalid") or parsed.get("Invalid") or ""
            #                 }
            #             elif isinstance(parsed, list):
            #                 note_list = parsed
            #                 note_map = {
            #                     "Color1": note_list[0] if len(note_list) > 0 else "",
            #                     "Color2": note_list[1] if len(note_list) > 1 else "",
            #                     "Color3": note_list[2] if len(note_list) > 2 else "",
            #                     "Invalid": note_list[3] if len(note_list) > 3 else ""
            #                 }
            #             else:
            #                 note_map = {"Color1": str(parsed), "Color2":"", "Color3":"", "Invalid":""}
            #         elif isinstance(notes_data, list):
            #             note_list = notes_data
            #             note_map = {
            #                 "Color1": note_list[0] if len(note_list) > 0 else "",
            #                 "Color2": note_list[1] if len(note_list) > 1 else "",
            #                 "Color3": note_list[2] if len(note_list) > 2 else "",
            #                 "Invalid": note_list[3] if len(note_list) > 3 else ""
            #             }
            #         else:
            #             # fallback: convert to string and put in Color1
            #             note_map = {"Color1": str(notes_data or ""), "Color2":"", "Color3":"", "Invalid":""}
            #     except Exception:
            #         note_map = {"Color1":"","Color2":"","Color3":"","Invalid":""}

            note_value = (note_map.get(color_key) or "").strip()

            if not note_value:
                continue

            # convert escaped newlines to actual newlines and set paragraphs (your existing logic)
            note_value = note_value.replace("\\n", "\n")
            lines = [line.strip() for line in note_value.split("\n") if line.strip()]

            shape.text_frame.text = ""
            shape.text_frame.word_wrap = True

            for i, line in enumerate(lines):
                p = shape.text_frame.paragraphs[0] if i == 0 else shape.text_frame.add_paragraph()
                p.text = line
                run = p.runs[0] if p.runs else p.add_run()
                font = run.font
                font.size = Pt(16)
                font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == "Prod_Test" and shape.has_text_frame:
            shape.height = Cm(3.42 if env_count >= 3 else 2.22)

        elif shape.name == "Dev_Value":
            if env_count < 3:
                slide.shapes._spTree.remove(shape._element)
        
        elif shape.name == "Dev_Text" and env_count == 2:
            slide.shapes._spTree.remove(shape._element)

        # --- Logic for Prod Tick/Cross ---
        if is_prod_red:
            if shape.name == "Prod_Value": # Hide tick if red
                shapes_to_remove.append(shape)
        else:
            if shape.name == "Prod_Value_Cross": # Hide cross if not red
                shapes_to_remove.append(shape)

        # --- Logic for Test Tick/Cross ---
        if is_test_red:
            if shape.name == "Test_Value": # Hide tick if red
                shapes_to_remove.append(shape)
        else:
            if shape.name == "Test_Value_Cross": # Hide cross if not red
                shapes_to_remove.append(shape)

    # Remove collected shapes AFTER iterating (safe removal)
    for sh in shapes_to_remove:
        try:
            slide.shapes._spTree.remove(sh._element)
        except Exception:
            # if removal fails for any reason, ignore (but optionally log)
            pass
    # Clear the list for safety if same var reused
    shapes_to_remove.clear()

    # ---Slide 4---
    slide4_data = data["slide4"]
    slide = prs.slides[3]
    headers = slide4_data["Storage_Utilization_Table"]["headers"]
    rows = slide4_data["Storage_Utilization_Table"]["rows"]
    used_val = float(rows[0][4])
    used_test = float(rows[1][4])
    rules = slide4_data["Colour_Rules"]
    
    if used_val >= rules["Color3"] or used_test >= rules["Color3"]:
        color_key = "Color3"
    elif used_val >= rules["Color2"] or used_test >= rules["Color2"]:
        color_key = "Color2"
    elif used_val >= rules["Color1"] or used_test >= rules["Color1"]:
        color_key = "Color1"
    else:
        color_key = "Invalid"
    circle_rgb = slide4_data["Circle_Color"][color_key]
    indicator_rgb = slide4_data["Indicator"][color_key]

    # Determine "red" status for Prod and Test individually
    is_prod_red = used_val >= rules["Color3"]
    is_test_red = used_test >= rules["Color3"]
    shapes_to_remove = []


    for shape in slide.shapes:
        if shape.name == "Storage_Utilization_Table" and shape.has_table:
            table = shape.table

            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
            
            for row_idx, row_data in enumerate(rows):
                for col_idx, cell_data in enumerate(row_data):
                    cell = table.cell(row_idx + 1, col_idx)
                    is_percent_col = col_idx in [len(row_data) - 1, len(row_data) - 2]
                    if isinstance(cell_data, (int, float)) and is_percent_col:
                         cell.text = f"{cell_data:.1f}%"
                    elif isinstance(cell_data, (int, float)):
                        cell.text = f"{cell_data:,}"
                    else:
                        cell.text = str(cell_data)
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(12)
                    p.alignment = PP_ALIGN.RIGHT
            
            total_rows_in_table = len(table.rows)
            needed_data_rows = len(rows)
            rows_to_delete = total_rows_in_table - 1 - needed_data_rows
            if rows_to_delete > 0:
                for _ in range(rows_to_delete):
                    delete_table_row(table, len(table.rows) - 1)
        
        elif shape.name == "Production_Storage_Usage_Chart" and shape.has_chart:
            chart = shape.chart
            chart.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
            chart_data = CategoryChartData()
            chart_data.categories = slide4_data["Production_Storage_Usage_Chart"]["Months"]
            for series_name, values in slide4_data["Production_Storage_Usage_Chart"].items():
                if series_name != "Months":
                    chart_data.add_series(series_name, values)
            chart.replace_data(chart_data)
        
            for s in chart.plots[0].series:
                if s.name == "Prod (GB)":
                    s.data_labels.number_format = '#,##0'
                    break
        
        elif shape.name == "Circle_Color":
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*circle_rgb)
        
        elif shape.name == "Indicator":
            fill = shape.fill
            if fill.type == 3:
                stops = fill.gradient_stops
                if len(stops) >= 2:
                    stops[0].color.rgb = RGBColor(*indicator_rgb)
                    stops[1].color.rgb = RGBColor(255, 255, 255)
        
        elif shape.name == "Notes_User_Input" and shape.has_text_frame:
                        # ---- Replace existing Notes_User_Input parsing with this ----
            notes_data = slide4_data.get("Notes_User_Input")

            # If DB returned a dict (jsonb -> Python dict), map it to the expected keys
            if isinstance(notes_data, dict):
                note_map = {
                    "Color1": notes_data.get("color1") or notes_data.get("Color1") or "",
                    "Color2": notes_data.get("color2") or notes_data.get("Color2") or "",
                    "Color3": notes_data.get("color3") or notes_data.get("Color3") or "",
                    "Invalid": notes_data.get("invalid") or notes_data.get("Invalid") or ""
                }
            #else:
                # existing/robust fallback: try to parse string/list as before
                # try:
                #     if isinstance(notes_data, str):
                #         parsed = json.loads(notes_data)
                #         if isinstance(parsed, dict):
                #             note_map = {
                #                 "Color1": parsed.get("color1") or parsed.get("Color1") or "",
                #                 "Color2": parsed.get("color2") or parsed.get("Color2") or "",
                #                 "Color3": parsed.get("color3") or parsed.get("Color3") or "",
                #                 "Invalid": parsed.get("invalid") or parsed.get("Invalid") or ""
                #             }
                #         elif isinstance(parsed, list):
                #             note_list = parsed
                #             note_map = {
                #                 "Color1": note_list[0] if len(note_list) > 0 else "",
                #                 "Color2": note_list[1] if len(note_list) > 1 else "",
                #                 "Color3": note_list[2] if len(note_list) > 2 else "",
                #                 "Invalid": note_list[3] if len(note_list) > 3 else ""
                #             }
                #         else:
                #             note_map = {"Color1": str(parsed), "Color2":"", "Color3":"", "Invalid":""}
                #     elif isinstance(notes_data, list):
                #         note_list = notes_data
                #         note_map = {
                #             "Color1": note_list[0] if len(note_list) > 0 else "",
                #             "Color2": note_list[1] if len(note_list) > 1 else "",
                #             "Color3": note_list[2] if len(note_list) > 2 else "",
                #             "Invalid": note_list[3] if len(note_list) > 3 else ""
                #         }
                #     else:
                #         # fallback: convert to string and put in Color1
                #         note_map = {"Color1": str(notes_data or ""), "Color2":"", "Color3":"", "Invalid":""}
                # except Exception:
                #     note_map = {"Color1":"","Color2":"","Color3":"","Invalid":""}

            note_value = (note_map.get(color_key) or "").strip()

            if not note_value:
                continue

            # convert escaped newlines to actual newlines and set paragraphs (your existing logic)
            note_value = note_value.replace("\\n", "\n")
            lines = [line.strip() for line in note_value.split("\n") if line.strip()]

            shape.text_frame.text = ""
            shape.text_frame.word_wrap = True

            for i, line in enumerate(lines):
                p = shape.text_frame.paragraphs[0] if i == 0 else shape.text_frame.add_paragraph()
                p.text = line
                run = p.runs[0] if p.runs else p.add_run()
                font = run.font
                font.size = Pt(16)
                font.color.rgb = RGBColor(0, 0, 0)

       
        elif shape.name == "Prod_Test" and shape.has_text_frame:
            shape.height = Cm(3.42 if env_count >= 3 else 2.22)
       
        elif shape.name == "Dev_Value":
            if env_count < 3:
                slide.shapes._spTree.remove(shape._element)
       
        elif shape.name == "Dev_Text" and env_count == 2:
            slide.shapes._spTree.remove(shape._element)

        # --- Logic for Prod Tick/Cross ---
        if is_prod_red:
            if shape.name == "Prod_Value": # Hide tick if red
                shapes_to_remove.append(shape)
        else:
            if shape.name == "Prod_Value_Cross": # Hide cross if not red
                shapes_to_remove.append(shape)

        # --- Logic for Test Tick/Cross ---
        if is_test_red:
            if shape.name == "Test_Value": # Hide tick if red
                shapes_to_remove.append(shape)
        else:
            if shape.name == "Test_Value_Cross": # Hide cross if not red
                shapes_to_remove.append(shape)

    # Remove collected shapes AFTER iterating (safe removal)
    for sh in shapes_to_remove:
        try:
            slide.shapes._spTree.remove(sh._element)
        except Exception:
            # if removal fails for any reason, ignore (but optionally log)
            pass
    # Clear the list for safety if same var reused
    shapes_to_remove.clear()

    # ---Slide 5---
    slide5_data = data["slide5"]
    slide = prs.slides[4]
    headers = slide5_data["Case_Status_Table"]["headers"]
    rows = slide5_data["Case_Status_Table"]["rows"]
    
    for shape in slide.shapes:
        if shape.name == "Case_Status_Table" and shape.has_table:
            table = shape.table

            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.bold = True

            for row_idx, row_data in enumerate(rows):
                for col_idx, cell_data in enumerate(row_data):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(cell_data)
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(12)

        elif shape.name == "Case_Trend_Chart" and shape.has_chart:
            chart = shape.chart
            chart_data = CategoryChartData()
            chart_data.categories = slide5_data["Case_Trend_Chart"]["Months"]
            for series_name, values in slide5_data["Case_Trend_Chart"].items():
                if series_name != "Months":
                    chart_data.add_series(series_name, values)
            chart.replace_data(chart_data)

        elif shape.name == "Open_Cases_Value" and shape.has_text_frame:
            p = shape.text_frame.paragraphs[0]
            p.text = str(slide5_data["Open_Cases_Value"])
            run = p.runs[0]
            font = run.font
            font.size = Pt(40)
            font.color.rgb = RGBColor(255, 255, 255)

    # ---Slide 7---
    slide7_data = data["slide7"]
    slide = prs.slides[6]
    for shape in slide.shapes:
        if shape.name == "Production_Availability_Chart" and shape.has_chart:
            chart = shape.chart
            availability = [val / 100 for val in slide7_data["Production_Availability_Chart"]["Availability"]]
            sla = [val / 100 for val in slide7_data["Production_Availability_Chart"]["SLA"]]
            chart_data = CategoryChartData()
            chart_data.categories = slide7_data["Production_Availability_Chart"]["Months"]
            chart_data.add_series("Availability", availability)
            chart_data.add_series("SLA", sla)
            chart.replace_data(chart_data)

            value_axis = chart.value_axis
            value_axis.minimum_scale = 0.93
            value_axis.maximum_scale = 1.0
            value_axis.tick_labels.number_format = '0.00%'

            for series in chart.series:
                if series.name == "Availability":
                    series.has_data_labels = True
                    series.data_labels.number_format = '0.00%'
                    series.data_labels.show_value = True
                else:
                    series.has_data_labels = False

        elif shape.name == "Production_User_Counts_Chart" and shape.has_chart:
            chart = shape.chart
            chart_data = CategoryChartData()
            # Remove Dev series if env_count is 2
            if env_count == 2:
                plot = chart.plots[0]
                series_list = list(plot.series)
                for idx, s in enumerate(series_list):
                    if s.name == "Dev":
                         # Remove via XML
                            ser_element = plot._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/chart}ser')[idx]
                            plot._element.remove(ser_element)
                            break

            chart_data.categories = slide7_data["Production_User_Counts_Chart"]["Months"]
            for series_name, values in slide7_data["Production_User_Counts_Chart"].items():
                if series_name != "Months":
                    chart_data.add_series(series_name, values)
            chart.replace_data(chart_data)

        elif shape.name == "Production_Storage_Usage_Chart" and shape.has_chart:
            chart = shape.chart
            chart.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
            chart_data = CategoryChartData()
            chart_data.categories = slide7_data["Production_Storage_Usage_Chart"]["Months"]
            for series_name, values in slide7_data["Production_Storage_Usage_Chart"].items():
                if series_name != "Months":
                    chart_data.add_series(series_name, values)
            chart.replace_data(chart_data)
            
            for s in chart.plots[0].series:
                if s.name == "Prod (GB)":
                    s.data_labels.number_format = '#,##0'
                    break

    prs.save(output_filename)
    print(f"Presentation saved as {output_filename}")
